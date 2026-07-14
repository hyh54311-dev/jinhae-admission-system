import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

FOLDER_ID = '1lxKMH8ssyeHSEv3rI2a24pD8qCT6FXo1'
SCOPES = ['https://www.googleapis.com/auth/drive']
BG_GYUWONGA = r"C:\Users\admin\.gemini\antigravity\brain\541b9e54-1667-4acf-816f-988959751015\bg_gyuwonga_1775181949007.png"
BG_INHYEON = r"C:\Users\admin\.gemini\antigravity\brain\541b9e54-1667-4acf-816f-988959751015\bg_inhyeon_1775181966145.png"

def add_custom_slide(prs, bg_img_path, title_text, bullet_points):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 1. 諛곌꼍 ?대?吏 ?쎌엯 (?꾩껜 李???린)
    slide.shapes.add_picture(bg_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # 2. ?쒕ぉ ?띿뒪??諛뺤뒪 ?쎌엯
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(1.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Malgun Gothic"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 230, 0) # 諛앹? ?몃???    
    # 援щ텇??(??媛뺤“)
    slide.shapes.add_shape(9, Inches(1), Inches(1.8), Inches(10), Inches(0.04))
    
    # 3. 蹂몃Ц ?띿뒪??諛뺤뒪 ?쎌엯
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(4.5))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    
    # 泥??⑤씫 泥섎━
    first_p = tf_body.paragraphs[0]
    first_pt = bullet_points[0]
    first_p.text = first_pt[0]
    first_p.level = first_pt[1]
    first_p.font.name = "Malgun Gothic"
    first_p.font.size = Pt(26 if first_pt[1] == 0 else 22)
    first_p.font.color.rgb = RGBColor(255, 100, 100) if first_pt[2] else RGBColor(255, 255, 255)
    
    # ?섎㉧吏 ?⑤씫 泥섎━
    for pt in bullet_points[1:]:
        text, indent_lvl, highlight = pt
        p = tf_body.add_paragraph()
        p.text = text
        p.level = indent_lvl
        p.line_spacing = 1.3
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(26 if indent_lvl == 0 else 22)
        # ?섏씠?쇱씠???щ????곕씪 遺됱???/ ?붿씠??        p.font.color.rgb = RGBColor(255, 120, 120) if highlight else RGBColor(255, 255, 255)
        if highlight:
            p.font.bold = True
            
def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ?щ씪?대뱶 1: ??댄?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(BG_GYUWONGA, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    t_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[?섎뒫?밴컯 臾명븰] 媛덈옒蹂듯빀 05"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "?덈궃?ㅽ뿄 <洹쒖썝媛> 諛??묒옄 誘몄긽 <?명쁽?뺥썑?? ?듭떖 ?대?"
    p2.font.name = "Malgun Gothic"
    p2.font.size = Pt(30)
    p2.font.color.rgb = RGBColor(255, 230, 0)
    p2.alignment = PP_ALIGN.CENTER
    
    # ?щ씪?대뱶 2: 洹쒖썝媛
    add_custom_slide(prs, BG_GYUWONGA, "1. ?덈궃?ㅽ뿄 <洹쒖썝媛> 媛쒓?", [
        ("二쇱젣: 遊됯굔 ?ы쉶 ???낆닔怨듬갑?섎뒗 ?ъ꽦???쒗깂怨??먮쭩", 0, False),
        ("媛덈옒: ?꾩쟾?섎뒗 媛???ㅻ옒??洹쒕갑 媛??, 0, False),
        ("?듭떖 ?ъ씤??1", 0, False),
        ("??怨쇨굅???꾨쫫?ㅼ?(?ㅻ퉰?붿븞) vs ?꾩옱???숆퀬 珥덈씪??硫대ぉ媛利?", 1, True),
        ("?듭떖 ?ъ씤??2", 0, False),
        ("??鍮꾪넻???⑦렪 ?먮쭩 ?띿뿉?쒕룄 ?꾩쓣 湲곕떎由щ뒗 洹뱀쟻 紐⑥닚??洹밸???, 1, True),
    ])
    
    # ?щ씪?대뱶 3: 洹쒖썝媛 以꾧굅由?1
    add_custom_slide(prs, BG_GYUWONGA, "2. <洹쒖썝媛> 怨쇨굅 ?뚯긽怨?諛⑺깢???⑦렪 ?먮쭩", [
        ("怨쇨굅????vs ?꾩옱????, 0, False),
        ("??15??利덉쓬?????덈????쇨뎬(?ㅻ퉰?붿븞)? ?대뵒濡?媛怨?, 1, False),
        ("??嫄곗슱 ?띿뿏 ?덈Т ?숆퀬 諛됱궡洹몃윭??紐⑥뒿(硫대ぉ媛利?肉?, 1, True),
        ("??遺洹臾대궓???섏쓽 湲곌뎄???붿옄, 議곕Ъ???쒓린?섎뒗援щ굹!", 1, False),
        ("?먮쭩?ㅻ윭???⑦렪???됱쟻", 0, False),
        ("???μ븞 ?좏삊 寃쎈컯???명깢??臾대━)? ?댁슱由щ뒗 ?⑦렪", 1, False),
        ("??諛깅쭏湲덊렪 李⑤┝?쇰줈 ?쇱쑀??湲곗깮吏? ?섍컙 ???뚯떇???딄?", 1, True),
    ])

    # ?щ씪?대뱶 4: 洹쒖썝媛 以꾧굅由?2
    add_custom_slide(prs, BG_GYUWONGA, "3. <洹쒖썝媛> ?낆닔怨듬갑 ?좊떖?붽낵 洹뱁븳??泥대뀗", [
        ("怨좎“?섎뒗 ?낆닔怨듬갑???몃줈?", 0, False),
        ("???⑥뼱吏??鍮쀬냼由? ?몄뼱???洹?쒕씪誘??ㅼ넄) ?뚮━", 1, False),
        ("??諛⑺솴怨??쒕쫫???щ옒?ㅺ퀬 嫄곕Ц怨??밴린湲?瑜??吏留?..", 1, False),
        ("??諛⑹? ??鍮꾧퀬 援ш끝媛꾩옣(?좏????꾪솄 援쎌씠 李쎌옄)留??딆뼱吏?, 1, True),
        ("泥대뀗 ?욎씤 ?띠쓽 留덈Т由?, 0, False),
        ("???쎌닔(嫄대꼸 ???녿뒗 臾? ?볦뿉 ?몄?議곗감 ?⑥젅??, 1, False),
        ("??李⑤씪由?二쎌뼱 泥쒕뀈 ?ㅼ쓽 ??蹂꾪븰)???섏뼱 ?뚯븘?ш퉴...", 1, True),
    ])

    # ?щ씪?대뱶 5: 以묒슂 媛쒕뀗
    add_custom_slide(prs, BG_GYUWONGA, "4. ?듭떖 ?섎뒫 媛쒕뀗! [媛먯젙?댁엯 vs 媛앷????곴?臾?", [
        ("媛먯젙?댁엯 (Empathy)", 0, False),
        ("????곸쓽 媛먯젙 = ?섏쓽 媛먯젙 (?묎컳???ы뵒)", 1, False),
        ("??<洹쒖썝媛> ??'洹?쒕씪誘??ㅼ넄)' (踰뚮젅媛 ?섏쿂???대떎怨?臾섏궗??", 1, True),
        ("媛앷????곴?臾?(Objective Correlative)", 0, False),
        ("??????먯껜???꾨Т ?앷컖 ?놁뼱????媛먯젙留??뗫낫?닿쾶 ?섎뒗 ?μ튂", 1, False),
        ("??(媛먯젙?댁엯??媛앷????곴?臾??덉뿉 ?ы븿?섎뒗 ??媛쒕뀗)", 1, False),
        ("??<洹쒖썝媛> ??'嫄곕Ц怨??밴린湲?' (洹몃깷 ?낃린吏留? ?곗＜?섎ŉ ?몃줈? 諛곌?)", 1, True),
    ])

    # ?щ씪?대뱶 6: ?명쁽?뺥썑??媛쒓?
    add_custom_slide(prs, BG_INHYEON, "5. ?묒옄 誘몄긽 <?명쁽?뺥썑?? ?밴컯", [
        ("議곗꽑 ?쒕? '3? 沅곸쨷 臾명븰'??諛깅?", 0, False),
        ("??湲곗궗?섍뎅 ~ 媛묒닠?섍뎅 ?쒓린 沅곸쨷 ?뷀닾瑜??ㅻ， ??궗/援?Ц ?뚯꽕", 1, False),
        ("?묓뭹???낆껜???뱀쭠", 0, False),
        ("??蹂듭옟???멸컙 ?대㈃蹂대떎??泥좎??섍쾶 ?좉낵 ?낆쓣 ??洹밸떒?쇰줈 ?뺤긽?뷀븿", 1, False),
        ("沅곴레?곸씤 吏묓븘 紐⑹쟻", 0, False),
        ("???좉탳??沅뚯꽑吏뺤븙 (?명쁽?뺥썑 蹂듭쐞???뺣떦???뺣낫, ?ν씗鍮?泥섎떒???뱀쐞??", 1, True),
    ])
    
    # ?щ씪?대뱶 7: ?몃Ъ ?由쎈룄
    add_custom_slide(prs, BG_INHYEON, "6. 洹밸떒???몃Ъ ?由?(??vs ??", [
        ("?명쁽?뺥썑 (珥덉??쇨? ?덈?????/ ?좉탳 ?뺣ぉ???붿떊)", 0, False),
        ("???꾧툑??已볦븘?닿퀬 ?띾컯?대룄 ?덈? ?먮쭩?섏? ?딆쓬 (洹쒖썝媛? ?ㅻ쫫)", 1, True),
        ("??援곗＜?????臾댄븳??遺?뺢낵 ?몃궡 (?듬떟???뺣룄???쒖쥌??", 1, False),
        ("?ν씗鍮?(???섏? ?먯슃怨???＜???곸쭠)", 0, False),
        ("??媛遺?μ젣媛 媛???먯삤?섎뒗 '?ш린(吏덊닾)'???꾩씠肄?, 1, True),
        ("?숈쥌 (媛遺?μ젣???덈? 沅뚮젰??", 0, False),
        ("??紐⑦븿???띿븘 ??젙???섎굹, 援?솗?닿린??媛뺣젰??鍮꾪뙋?먯꽌???댁쭩 鍮쀪꺼媛?, 1, False),
    ])

    # ?щ씪?대뱶 8: 以꾧굅由??붿빟
    add_custom_slide(prs, BG_INHYEON, "7. 癒몃┸?띿뿉 諛붾줈 苑귦엳???ш굔???꾧컻", [
        ("諛쒕떒: 15?몄쓽 ?대┛ ?명쁽?뺥썑媛 ?숈쥌????踰덉㎏ 遺?몄씠 ?섎떎", 0, False),
        ("?꾧컻: ?섏?留??꾨뱾???녹? 嫄??ν씗鍮? 臾댁옄鍮꾪븳 ?뚰빐媛 ?쒖옉??, 0, False),
        ("?꾧린: ?ν씗鍮?移섎쭧??뿉 ?띿? ?숈쥌, ?명쁽?뺥썑瑜??쒖씤?쇰줈 ?먯쐞", 0, True),
        ("??已볤꺼???덇뎅???ш??먯꽌???명쁽?뺥썑???쒖쥌?쇨? 移⑥갑???좎?", 1, False),
        ("?덉젙: ?ν씗鍮덉씠 以묒쟾???섏뿀?쇰굹 洹??⑥븙吏덉뿉 ?숈쥌??吏덈젮踰꾨┝", 0, True),
        ("寃곕쭚: ?숈쥌 ??ㅺ컖?????명쁽?뺥썑 蹂듭쐞 ???ν씗鍮?李멸탳???ъ빟)", 0, False),
    ])

    # ?щ씪?대뱶 9: ?쒖닠?먯쓽 媛쒖엯
    add_custom_slide(prs, BG_INHYEON, "8. 怨좎쟾 ?뚯꽕 ?④낏 臾몄젣! [?쒖닠?먯쓽 媛쒖엯]", [
        ("?쒖닠?먯쓽 媛쒖엯 (?몄쭛?먯쟻 ?쇳룊)", 0, False),
        ("??愿李곗옄?ъ빞 ???뚯쫰 ???쒖닠?먭? 遺덉뫁 ?섑????곹솴???덉닔?먭퀬 ?됯???, 1, False),
        ("蹂몃Ц ?덉떆: ?쒖닠?먯쓽 ?몄븷 ?곗쭚", 0, False),
        ("??\"?댁컡 媛?⑥튂 ?꾨땲?섎━?? 留뚮Ъ???ы띁?섍퀬 ?섎뒛??媛먮룞?덈떎\"", 1, True),
        ("?대젃寃???볤퀬 ?쇳룊???좊━???④낵?", 0, False),
        ("???낆옄?ㅼ쓣 媛먯젙?곸쑝濡??몃뇤(?)?섏뿬 ?먯뿰?ㅻ젅 ?명쁽?뺥썑 ?몄쓣 ?ㅺ쾶 ??, 1, True),
        ("???ν씗鍮덉쓽 ?낇뻾?????遺꾨끂 ?섏튂瑜?理쒕?濡??뚯뼱?щ┝", 1, False),
    ])

    prs.save("worksheet_presentation_v2.pptx")
    print("Local V2 presentation saved as worksheet_presentation_v2.pptx")

def upload_presentation():
    creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    service = build('drive', 'v3', credentials=creds)
    
    file_metadata = {
        'name': '[?щ씪?대뱶 V2] 2027 ?섎뒫?밴컯 臾명븰 媛?낆꽦 洹밸???,
        'parents': [FOLDER_ID],
        'mimeType': 'application/vnd.google-apps.presentation'
    }
    
    media = MediaFileUpload('worksheet_presentation_v2.pptx',
                            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                            resumable=True)
                            
    print("Uploading to Google Drive and converting to Slides...")
    file = service.files().create(body=file_metadata, media_body=media, fields='id, webViewLink').execute()
    print(f"Presentation uploaded! Link: {file.get('webViewLink')}")

if __name__ == "__main__":
    create_presentation()
    upload_presentation()
