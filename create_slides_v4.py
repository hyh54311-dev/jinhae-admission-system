import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from PIL import Image

FOLDER_ID = '1lxKMH8ssyeHSEv3rI2a24pD8qCT6FXo1'
SCOPES = ['https://www.googleapis.com/auth/drive']
BG_GYUWONGA_RAW = r"C:\Users\admin\.gemini\antigravity\brain\541b9e54-1667-4acf-816f-988959751015\bg_gyuwonga_1775181949007.png"
BG_INHYEON_RAW = r"C:\Users\admin\.gemini\antigravity\brain\541b9e54-1667-4acf-816f-988959751015\bg_inhyeon_1775181966145.png"

BG_GYUWONGA = r"C:\Users\admin\.gemini\antigravity\brain\541b9e54-1667-4acf-816f-988959751015\bg_gyuwonga_light.png"
BG_INHYEON = r"C:\Users\admin\.gemini\antigravity\brain\541b9e54-1667-4acf-816f-988959751015\bg_inhyeon_light.png"

TITLE_FONT = "Gmarket Sans Medium" 
BODY_FONT = "Noto Sans KR"

# 諛앹? ???붿씠?? ?띿뒪??媛뺤“ 而щ윭 ?ㅼ젙
COLOR_TEXT_MAIN = RGBColor(30, 30, 30)       # 吏숈? ?뚭???COLOR_TITLE = RGBColor(0, 51, 102)           # 源딄퀬 吏꾪븳 ?ㅼ씠鍮?(?쒕ぉ)
COLOR_HIGHLIGHT = RGBColor(190, 0, 0)        # ?꾩＜ 吏꾪븳 鍮④컙??(媛뺤“)

def wash_image(src, dst):
    try:
        img = Image.open(src).convert("RGB")
        white = Image.new('RGB', img.size, 'white')
        # ?곗깋??70% ?뺣룄 ?욎뼱 ?섏콈?붿쿂???꾩＜ ?고븳 ?뚰꽣留덊겕 諛곌꼍?쇰줈 留뚮벀
        washed = Image.blend(img, white, 0.7) 
        washed.save(dst)
    except Exception as e:
        print(f"?대?吏 釉붾젋???ㅽ뙣: {e}")

def add_custom_slide(prs, bg_img_path, title_text, bullet_points):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    slide.shapes.add_picture(bg_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(1.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = TITLE_FONT
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TITLE
    
    slide.shapes.add_shape(9, Inches(1), Inches(1.8), Inches(10), Inches(0.04))
    
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(4.5))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    
    first_p = tf_body.paragraphs[0]
    first_pt = bullet_points[0]
    first_p.text = first_pt[0]
    first_p.level = first_pt[1]
    first_p.font.name = BODY_FONT
    first_p.font.size = Pt(26 if first_pt[1] == 0 else 22)
    first_p.font.color.rgb = COLOR_HIGHLIGHT if first_pt[2] else COLOR_TEXT_MAIN
    
    for pt in bullet_points[1:]:
        text, indent_lvl, highlight = pt
        p = tf_body.add_paragraph()
        p.text = text
        p.level = indent_lvl
        p.line_spacing = 1.3
        p.font.name = BODY_FONT
        p.font.size = Pt(26 if indent_lvl == 0 else 22)
        p.font.color.rgb = COLOR_HIGHLIGHT if highlight else COLOR_TEXT_MAIN
        if highlight:
            p.font.bold = True
            
def create_presentation():
    # 諛곌꼍 ?대?吏瑜??곗깋?쇰줈 ?뚯떆(Wash)?섏뿬 ?щ챸?꾨? ?믪엫
    wash_image(BG_GYUWONGA_RAW, BG_GYUWONGA)
    wash_image(BG_INHYEON_RAW, BG_INHYEON)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(BG_GYUWONGA, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    t_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[?섎뒫?밴컯 臾명븰] 媛덈옒蹂듯빀 05"
    p.font.name = TITLE_FONT
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "?덈궃?ㅽ뿄 <洹쒖썝媛> 諛??묒옄 誘몄긽 <?명쁽?뺥썑?? ?듭떖 ?대?"
    p2.font.name = TITLE_FONT
    p2.font.size = Pt(30)
    p2.font.color.rgb = COLOR_TITLE
    p2.alignment = PP_ALIGN.CENTER
    
    add_custom_slide(prs, BG_GYUWONGA, "1. ?덈궃?ㅽ뿄 <洹쒖썝媛> 媛쒓?", [
        ("二쇱젣: 遊됯굔 ?ы쉶 ???낆닔怨듬갑?섎뒗 ?ъ꽦???쒗깂怨??먮쭩", 0, False),
        ("媛덈옒: ?꾩쟾?섎뒗 媛???ㅻ옒??洹쒕갑 媛??, 0, False),
        ("?듭떖 ?ъ씤??1", 0, False),
        ("??怨쇨굅???덈????쇨뎬(?ㅻ퉰?붿븞) vs ?꾩옱???숆퀬 珥덈씪??硫대ぉ媛利?", 1, True),
        ("?듭떖 ?ъ씤??2", 0, False),
        ("??鍮꾪넻???⑦렪 ?먮쭩 ?띿뿉?쒕룄 ?꾩쓣 湲곕떎由щ뒗 ?댁쨷???쒕룄??洹밸???, 1, True),
    ])
    
    add_custom_slide(prs, BG_GYUWONGA, "2. <洹쒖썝媛> 湲??? 鍮쏅굹??怨쇨굅? 諛⑺깢???꾩옱 ?⑦렪", [
        ("怨쇨굅????vs ?꾩옱?????щ튆 ?鍮?, 0, False),
        ("??15??利덉쓬?????덈????쇨뎬(?ㅻ퉰?붿븞)? ?꾨?泥??대뵒濡?媛怨?, 1, False),
        ("??嫄곗슱 ?띿뿏 ?덈Т ?숆퀬 諛됱궡洹몃윭??紐⑥뒿(硫대ぉ媛利?肉?, 1, True),
        ("??遺洹臾대궓???섏쓽 湲곌뎄???붿옄, 議곕Ъ???쒓린?섎뒗援щ굹!", 1, False),
        ("?먮쭩?ㅻ윭???⑦렪???꾩옱 沅ㅼ쟻", 0, False),
        ("???μ븞 ?좏삊 寃쎈컯??二쇱깋?↔린瑜?利먭린???명깢??臾대━)? ?댁슱由щ뒗 ?⑦렪", 1, False),
        ("????湲곗깮吏??쇱쑀?????섍컙 ???????듭쓽 ?뚯떇議곗감 ?딄?", 1, True),
    ])

    add_custom_slide(prs, BG_GYUWONGA, "3. <洹쒖썝媛> ??寃? ?낆닔怨듬갑 ?몃줈?怨?洹뱁븳??泥대뀗", [
        ("怨좎“?섎뒗 ?낆닔怨듬갑???몃줈???뚮━??, 0, False),
        ("???⑥뼱吏??鍮쀬냼由? 踰쎌삤?숈뿉 ?몄뼱???洹?쒕씪誘??ㅼ넄) ?뚮━", 1, False),
        ("??諛⑺솴怨??쒕쫫???щ옒?ㅺ퀬 嫄곕Ц怨??밴린湲?瑜??吏留?..", 1, False),
        ("??諛⑹? ??鍮꾧퀬 援ш끝媛꾩옣(?좏????꾪솄 援쎌씠 李쎌옄)留??딆뼱吏???꾪뵒", 1, True),
        ("泥대뀗 ?욎씤 ?띠쓽 ?몄벝??留덈Т由?, 0, False),
        ("???쎌닔(嫄대꼸 ???녿뒗 臾? 媛숈? ?μ븷臾??볦뿉 ?몄?議곗감 ?⑥젅??, 1, False),
        ("??轅덉뿉?쒖“李??⑦렪??紐?蹂대뒓?? 李⑤씪由?二쎌뼱 泥쒕뀈 ???숈씠 ?좉퉴...", 1, True),
    ])

    add_custom_slide(prs, BG_GYUWONGA, "4. ?듭떖 ?섎뒫 媛쒕뀗! [媛먯젙?댁엯 vs 媛앷????곴?臾?", [
        ("媛먯젙?댁엯 (Empathy)", 0, False),
        ("????곸쓽 媛먯젙 = ?섏쓽 媛먯젙 (?덈룄 ?섏쿂???ы봽援щ굹!)", 1, False),
        ("??<洹쒖썝媛> ??'洹?쒕씪誘??ㅼ넄)' (踰뚮젅媛 ?섏쿂??'?대떎'怨??쒗쁽??", 1, True),
        ("媛앷????곴?臾?(Objective Correlative)", 0, False),
        ("???곗슱????鍮꾧? ?ㅻ㈃ ???몃∼?? ??媛먯젙???뗫낫?닿쾶 ?섎뒗 ?꾧뎄", 1, False),
        ("??(媛먯젙?댁엯??媛앷????곴?臾?移댄뀒怨좊━ ?덉뿉 ?ы븿??", 1, False),
        ("??<洹쒖썝媛> ??'嫄곕Ц怨??밴린湲?' (?닿? 洹몃깷 移섎뒗 ?낃린??肉먯씠吏留??몃줈? 諛곌?)", 1, True),
    ])

    add_custom_slide(prs, BG_INHYEON, "5. ?묒옄 誘몄긽 <?명쁽?뺥썑?? ?밴컯", [
        ("議곗꽑 ?쒕? ?ν뻾 ?뚰뭾 '3? 沅곸쨷 臾명븰'??諛깅?", 0, False),
        ("??湲곗궗?섍뎅 ~ 媛묒닠?섍뎅 ?쒓린 沅곸쨷 ?뷀닾瑜??ㅻ， ??궗 ?뚯꽕", 1, False),
        ("???뚯꽕 ???몃Ъ?ㅼ쓽 ?뱀쭠??", 0, False),
        ("??蹂듭옟???멸컙???ш꼍 蹂?붾낫?ㅻ뒗 泥좎??섍쾶 ?좉낵 ?낆쓣 ??洹밸떒?쇰줈 ?뺤긽?뷀븿", 1, False),
        ("沅곴레?곸씤 吏묓븘 紐⑹쟻", 0, False),
        ("???좉탳??沅뚯꽑吏뺤븙 ?ъ긽???뱁샇", 1, True),
        ("???명쁽?뺥썑 蹂듭쐞???뺣떦???뺣낫 諛??ν씗鍮??낇뻾 泥섎떒???뱀쐞???낆쬆", 1, True),
    ])
    
    add_custom_slide(prs, BG_INHYEON, "6. 洹밸떒???몃Ъ ?由?(??vs ??", [
        ("?명쁽?뺥썑 (珥덉??쇨? ?덈?????/ ?좉탳 ?뺣ぉ???붿떊)", 0, False),
        ("???꾧툑??已볦븘?닿퀬 ?ш?濡??댁퀜???덈? ?먮쭩?섏? ?딆쓬", 1, True),
        ("??援곗＜?????臾댄븳??遺?뺢낵 ?몃궡 (?듬떟???뺣룄??洹뱁븳???쒖쥌??", 1, False),
        ("?ν씗鍮?(???섏? ?먯슃怨???＜???곸쭠)", 0, False),
        ("??媛遺?μ젣媛 媛???먯삤?섎뒗 ?ㅻ쭔遺덉넀怨?'?ш린(吏덊닾)'???꾩씠肄?, 1, True),
        ("?숈쥌 (媛遺?μ젣???덈? 沅뚮젰??", 0, False),
        ("??紐⑦븿???띿븘 ??젙???섎굹, 援?솗?닿린??媛뺣젰??鍮꾪뙋?먯꽌???댁쭩 鍮쀪꺼媛?, 1, False),
    ])

    add_custom_slide(prs, BG_INHYEON, "7. 癒몃┸?띿뿉 諛붾줈 苑귦엳???ш굔???꾧컻留?, [
        ("諛쒕떒: 15?몄쓽 ?대┛ ?명쁽?뺥썑媛 ?숈쥌????踰덉㎏ 遺?몄씠 ?섎떎", 0, False),
        ("?꾧컻: ?섏?留??꾨뱾???녹? 嫄??ν씗鍮? 臾댁옄鍮꾪븳 紐⑦븿???쒖옉??, 0, False),
        ("?꾧린: ?ν씗鍮?移섎쭧??뿉 ?띿? ?숈쥌, ?명쁽?뺥썑瑜??쒖씤?쇰줈 ?먯쐞", 0, True),
        ("??已볤꺼???덇뎅???ш??먯꽌???명쁽?뺥썑???쒖쥌?쇨? 移⑥갑?④낵 ?몃궡 ?좎?", 1, False),
        ("?덉젙: ?ν씗鍮덉씠 ?쒖빳 以묒쟾???섏뿀?쇰굹 洹??⑥븙吏덉뿉 ?숈쥌??吏덈젮踰꾨┝", 0, True),
        ("寃곕쭚: ?숈쥌????ㅺ컖?????명쁽?뺥썑 ?붾젮??蹂듭쐞 ???ν씗鍮?李멸탳???ъ빟)", 0, False),
    ])

    add_custom_slide(prs, BG_INHYEON, "8. 怨좎쟾 ?뚯꽕 ?④낏 臾몄젣! [?쒖닠?먯쓽 媛쒖엯]", [
        ("?쒖닠?먯쓽 媛쒖엯 (?몄쭛?먯쟻 ?쇳룊) ?대??", 0, False),
        ("???щ챸?멸컙?대뜕 移대찓?쇰㎤(?쒖닠????遺덉뫁 ?섑????곹솴???덉닔?먭퀬 ?됯???, 1, False),
        ("蹂몃Ц ?덉떆: ?쒖닠?먯쓽 媛먯젙 ?댁엯??鍮??곗쭚", 0, False),
        ("??\"?댁컡 媛?⑥튂 ?꾨땲?섎━?? 留뚮Ъ???ы띁?섍퀬 ?섎뒛??媛먮룞?덈떎\"", 1, True),
        ("?대젃寃???볤퀬 ?쇳룊???좊━硫??앷린???④낵?", 0, False),
        ("???낆옄?ㅼ뿉寃??먯뿰?ㅻ젅 ?명쁽?뺥썑 ?몄쓣 ?ㅺ쾶 ?섎뒗 媛뺣젰??媛?ㅻ씪?댄똿(?) ?μ튂", 1, True),
        ("???ν씗鍮덉쓽 ?낇뻾?????遺꾨끂 ?섏튂瑜??낆옄? ?④퍡 理쒕?濡??뚯뼱?щ┝", 1, False),
    ])

    prs.save("worksheet_presentation_v4.pptx")
    print("Local V4 presentation saved as worksheet_presentation_v4.pptx")

def upload_presentation():
    creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    service = build('drive', 'v3', credentials=creds)
    
    file_metadata = {
        'name': '[?щ씪?대뱶 V4] 2027 ?섎뒫?밴컯 臾명븰 (?붿씠???띿뒪爾?理쒖쟻??',
        'parents': [FOLDER_ID],
        'mimeType': 'application/vnd.google-apps.presentation'
    }
    
    media = MediaFileUpload('worksheet_presentation_v4.pptx',
                            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                            resumable=True)
                            
    print("Uploading to Google Drive and converting to Slides...")
    file = service.files().create(body=file_metadata, media_body=media, fields='id, webViewLink').execute()
    print(f"Presentation uploaded! Link: {file.get('webViewLink')}")

if __name__ == "__main__":
    create_presentation()
    upload_presentation()
