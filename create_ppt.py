from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ?뚰겕?뚮줈 ?붿옄??洹쒖튃???곕Ⅸ 湲?먯깋
    color_title = RGBColor(0, 51, 102)
    color_body = RGBColor(30, 30, 30)
    color_hl = RGBColor(190, 0, 0)
    
    def add_slide(title_text, body_bullets):
        # 6踰??덉씠?꾩썐: 鍮??붾㈃
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # ?쒕ぉ ?띿뒪?몃컯??        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(50)
        run.font.bold = True
        run.font.color.rgb = color_title

        # 蹂몃Ц ?띿뒪?몃컯??        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(4.5))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        
        for i, text in enumerate(body_bullets):
            if i == 0:
                p = tf_body.paragraphs[0]
            else:
                p = tf_body.add_paragraph()
            
            p.level = 0
            p.space_after = Pt(25)
            run = p.add_run()
            run.text = "??" + text
            run.font.name = "Malgun Gothic"
            run.font.size = Pt(36)
            run.font.bold = True
            
            # ?뚭뎅, 鍮꾧레, ?ㅻ컻?????듭떖 ?ㅼ썙???ы븿 ??媛뺤“ 而щ윭
            if any(keyword in text for keyword in ["?뚭뎅", "遺뺢눼", "?ㅻ컻??, "?꾨씫", "?щ쭩"]):
                run.font.color.rgb = color_hl
            else:
                run.font.color.rgb = color_body
                
        return slide

    # 1. 硫붿씤 ??댄? ?щ씪?대뱶 (媛????湲??
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.333), Inches(2.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[臾명븰] '?ㅻ컻?? ?꾩껜 以꾧굅由??붿빟"
    r.font.name = "Malgun Gothic"
    r.font.bold = True
    r.font.size = Pt(65)
    r.font.color.rgb = color_title
    
    # 2. 蹂몃Ц ?щ씪?대뱶 援ъ꽦 (?꾩엯, 諛쒕떒, ?꾧린, ?덉젙, 寃곕쭚)
    slides_data = [
        ("1. ?꾩엯 (?쒕???諛곌꼍 ?명똿)", [
            "?쒓컙??諛곌꼍: 6.25 ?꾩웳 吏곹썑 1950?꾨?", 
            "怨듦컙??諛곌꼍: ?덈쭩?곸씠怨?媛?쒗븳 ?쒖슱 ?대갑珥?, 
            "遺꾩쐞湲? ?꾩웳???곹쓷怨?洹뱀떖??鍮꾧레???쒕?"
        ]),
        ("2. 諛쒕떒 (臾몄젣??吏뺤“)", [
            "二쇱씤怨?泥좏샇??泥섏젅?섍퀬 湲곌뎄??媛議깆궗", 
            "'媛??'留??몄튂??誘몄튇 ?대㉧?덉? 媛?쒗븳 ?대┝", 
            "遺?묒쓽 臾닿쾶??吏볥닃由?泥좏샇???ы빐吏??'移섑넻'"
        ]),
        ("3. ?꾧린 (媛덈벑??怨좎“)", [
            "?숈깮 ?곹샇: ?앹〈???꾪빐 ?묒떖??踰꾨━湲곕줈 寃곗떖??, 
            "?숈깮 紐낆닕: 媛議깆쓣 遺?묓븯??誘멸뎔 ?꾩븞遺濡??꾨씫??, 
            "泥좏샇??怨좊뇤: ?묒떖 ?뚮Ц??媛議깆쓣 援ы븯吏 紐삵븿"
        ]),
        ("4. ?덉젙 (鍮꾧레????컻)", [
            "?앹〈???꾪빐 媛뺣룄吏볦쓣 ???곹샇媛 寃쎌같??援ъ냽??, 
            "留뚯궘?대뜕 ?꾨궡??媛묒옉?ㅻ윭???щ쭩 ?뚯떇", 
            "紐⑤뱺 ?щ쭩??臾대꼫???대━??泥좏샇??硫섑깉 遺뺢눼"
        ]),
        ("5. 寃곕쭚 (?ㅻ컻?꾩쓽 ?섎?)", [
            "吏湲뗭?湲뗮븳 移섑넻(?볥뒗 ????寃곌뎅 紐⑤몢 ??戮묒븘??, 
            "紐⑹쟻吏瑜??껉퀬 '媛??' ?섏껌???⑹떥???앹떆瑜???, 
            "諛⑺뼢???껋? ?멸컙, ?쇰챸 '?ㅻ컻?? 媛숈? 議댁옱??鍮꾧레"
        ])
    ]
    
    for title, bullets in slides_data:
        add_slide(title, bullets)
        
    save_path = "D:\\OneDrive - 寃쎌긽?⑤룄援먯쑁泥?\諛뷀깢 ?붾㈃\\2?숇뀈臾명븰_?ㅻ컻??以꾧굅由??붿빟.pptx"
    prs.save(save_path)
    print("??μ셿猷? " + save_path)

if __name__ == '__main__':
    create_ppt()
