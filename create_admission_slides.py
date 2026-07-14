import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- 디자인 테마 설정 ---
TITLE_FONT = "Gmarket Sans Medium" 
BODY_FONT = "Noto Sans KR"

# 네이비 & 화이트 테마
COLOR_BG_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_MAIN = RGBColor(40, 40, 40)        # 짙은 회색
COLOR_TITLE = RGBColor(12, 35, 64)            # 진해고 딥 네이비
COLOR_HIGHLIGHT = RGBColor(211, 47, 47)       # 짙은 레드 (강조용)
COLOR_MUTED = RGBColor(120, 120, 120)         # 연한 회색

LOGO_PATH = "school_logo.png"

def add_header_and_logo(slide, prs):
    """모든 일반 슬라이드에 적용되는 공통 헤더 라인과 우측 상단 로고"""
    # 우측 상단 로고 이미지 (school_logo.png가 존재하는 경우 삽입)
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, Inches(12.0), Inches(0.3), width=Inches(1.0))
        except Exception as e:
            print(f"로고 삽입 실패: {e}")

    # 상단 장식 라인 (딥 네이비)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_TITLE
    shape.line.fill.background()

def add_title(slide, text):
    """슬라이드 타이틀 텍스트 추가"""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.0), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = TITLE_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    return title_box

def create_admission_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: 표지 (Title Slide)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # 배경 데코 박스 (네이비 블록)
    bg_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLOR_TITLE
    bg_shape.line.fill.background()
    
    # 메인 타이틀
    title_box = slide1.shapes.add_textbox(Inches(5.0), Inches(2.2), Inches(7.5), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_sub = tf.paragraphs[0]
    p_sub.text = "2026학년도 신입생 유치를 위한"
    p_sub.font.name = BODY_FONT
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = COLOR_MUTED
    p_sub.space_after = Pt(10)
    
    p_main = tf.add_paragraph()
    p_main.text = "진해고등학교 입학 설명회"
    p_main.font.name = TITLE_FONT
    p_main.font.size = Pt(48)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_TITLE
    p_main.space_after = Pt(10)
    
    p_desc = tf.add_paragraph()
    p_desc.text = "변화하는 대입 제도, 왜 진해고가 정답인가?"
    p_desc.font.name = BODY_FONT
    p_desc.font.size = Pt(18)
    p_desc.font.color.rgb = COLOR_HIGHLIGHT
    p_desc.font.bold = True
    
    # 학교 로고 (표지 중앙 네이비 블록 안 배치)
    if os.path.exists(LOGO_PATH):
        try:
            slide1.shapes.add_picture(LOGO_PATH, Inches(1.2), Inches(2.7), width=Inches(2.1))
        except Exception as e:
            print(f"표지 로고 삽입 실패: {e}")

    # ==========================================
    # SLIDE 2: 찾아가는 설명회 - 브리핑 개요
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, "찾아가는 설명회: 부장교사 Professional 브리핑")
    add_header_and_logo(slide2, prs)
    
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    pts = [
        ("🎯 대형 강당 브리핑의 목표", 0, True),
        ("- 교실별 분산 진행 대신 대형 강당 진행으로 학교의 규모감과 무게감 선사", 1, False),
        ("- 중3 학생들에게 '왜 지금 진해고를 선택해야 하는가'에 대한 확신 부여", 1, False),
        ("📢 핵심 메시지 설계", 0, True),
        ("- '5등급제 도입(고교학점제)'에 따른 대입 지형 변화 선제 분석", 1, False),
        ("- 학생 수가 많을 때 보장되는 선택과목 개설 다양성과 내신 안전망 강조", 1, False),
        ("💡 운영 전략: 투트랙(Two-Track) 시너지", 0, True),
        ("- 부장교사의 Professional 브리핑 (거시 비전, 신뢰도 확보) [15분]", 1, False),
        ("- 재학생 도우미의 솔직 담백 Q&A 및 토크쇼 (감성 터치, 공감대 형성) [15분]", 1, False)
    ]
    
    for idx, (text, level, highlight) in enumerate(pts):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        p.text = text
        p.level = level
        p.line_spacing = 1.3
        p.space_after = Pt(8)
        p.font.name = BODY_FONT
        p.font.size = Pt(24 if level == 0 else 20)
        p.font.color.rgb = COLOR_TITLE if level == 0 else COLOR_TEXT_MAIN
        if highlight:
            p.font.bold = True

    # ==========================================
    # SLIDE 3: 찾아가는 설명회 - Q&A 큐카드 시나리오
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "재학생 입학홍보단 '솔직 담백' Q&A 큐카드")
    add_header_and_logo(slide3, prs)
    
    # 3개의 큐카드를 가로 배치하는 카드 레이아웃
    card_width = Inches(3.6)
    card_height = Inches(4.2)
    card_y = Inches(2.2)
    card_xs = [Inches(0.8), Inches(4.8), Inches(8.8)]
    
    qna_data = [
        ("Q1. 매점과 급식", "진해고 급식 양도 푸짐하고 매월 특식 퀄리티 최상!\n2교시 쉬는 시간 매점 닭강정과 만두 쟁탈전은 진해고의 활력소!", COLOR_TITLE),
        ("Q2. 현실적인 내신 걱정", "인원 적은 소규모 학교는 실수 한 번에 등급 폭망..\n진해고는 인원이 많아 실수 만회가 가능한 든든한 '안전망' 존재!", COLOR_HIGHLIGHT),
        ("Q3. 남고의 끈끈함 분위기", "칙칙하고 거칠다는 소문은 NO!\n체육대회/축제 우주 최강 단합력, 평생 갈 진짜 친구 사귀는 곳!", COLOR_TITLE)
    ]
    
    for i, (q_title, q_body, color) in enumerate(qna_data):
        # 외곽 카드 상자
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_xs[i], card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 텍스트 배치
        tb = slide3.shapes.add_textbox(card_xs[i] + Inches(0.1), card_y + Inches(0.1), card_width - Inches(0.2), card_height - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_q = tf.paragraphs[0]
        p_q.text = q_title
        p_q.font.name = TITLE_FONT
        p_q.font.size = Pt(22)
        p_q.font.bold = True
        p_q.font.color.rgb = color
        p_q.space_after = Pt(14)
        
        p_a = tf.add_paragraph()
        p_a.text = q_body
        p_a.font.name = BODY_FONT
        p_a.font.size = Pt(17)
        p_a.line_spacing = 1.3
        p_a.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 4: 교내 설명회 - 대규모 남고의 구조적 강점
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "진해고등학교만이 가지는 독보적 구조적 강점")
    add_header_and_logo(slide4, prs)
    
    content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf4 = content_box.text_frame
    tf4.word_wrap = True
    
    pts = [
        ("🚫 여학생들과의 내신 경쟁 스트레스 원천 차단", 0, True),
        ("- 공학 고교 대비 꼼꼼한 수행평가 등 남학생이 상대적으로 겪을 수 있는 내신 불리함 완벽 방지", 1, False),
        ("- 오직 본인의 성취도와 역량에만 집중하여 공정하고 안정적으로 내신을 확보하는 남학교 환경", 1, False),
        ("🎯 교내 진로 선택과목 100% 개설 (폐강/통폐합 우려 Zero)", 0, True),
        ("- 학생부종합전형의 핵심: 본인의 진로에 맞는 기하, 물리학Ⅱ, 프로그래밍, 인공지능기초 등 심화 이수 필수", 1, False),
        ("- 소규모 고교와 달리 학년별 약 300명의 두터운 학생 수 덕분에 모든 진로 선택과목이 교내 정규 과정으로 개설", 1, False),
        ("- 다른 학교와 공동교육과정(위탁)으로 시간과 동선을 낭비할 필요 없이 학교 안에서 완벽히 해결 가능", 1, True)
    ]
    
    for idx, (text, level, highlight) in enumerate(pts):
        p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
        p.text = text
        p.level = level
        p.line_spacing = 1.3
        p.space_after = Pt(8)
        p.font.name = BODY_FONT
        p.font.size = Pt(24 if level == 0 else 20)
        p.font.color.rgb = COLOR_TITLE if level == 0 or not highlight else COLOR_HIGHLIGHT
        if highlight:
            p.font.bold = True

    # ==========================================
    # SLIDE 5: 교내 설명회 - 다수 인원의 강점 (표 포함)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_title(slide5, "선택과목 이수 및 내신 등급의 치트키: '인원수'")
    add_header_and_logo(slide5, prs)
    
    # 상단 멘트
    desc_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "※ 고교학점제 5등급제(1등급 상위 10%) 기준 수강 인원별 등급 문호 비교"
    p_desc.font.name = BODY_FONT
    p_desc.font.size = Pt(18)
    p_desc.font.color.rgb = COLOR_HIGHLIGHT
    p_desc.font.bold = True
    
    # 표(Table) 생성: 4행 4열
    rows, cols = 4, 4
    left, top, width, height = Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.5)
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 열 너비 조절
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(3.2)
    
    # 표 헤더 채우기
    headers = ["구분 (과목별 수강 인원)", "1등급 인원 (상위 10%)", "2등급 누적 (상위 34%)", "학부모 안심 포인트"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = TITLE_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_BG_WHITE
        p.font.bold = True
        
    # 표 데이터 채우기
    row_data = [
        ["소규모 학교/과목 (20명 수강)", "단 2명", "누적 6명", "실수 한 번에 등급 급락, 압박감 극대화"],
        ["중규모 학교/과목 (40명 수강)", "4명", "누적 13명", "내신 경쟁 과열로 인한 학업 스트레스"],
        ["진해고 우위 과목 (100명 수강)", "10명 (5배 확보)", "누적 34명", "내신 획득의 문호가 넓어 시험 만회 가능"]
    ]
    
    for row_idx, data in enumerate(row_data):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            # 마지막 행 강조
            if row_idx == 2:
                cell.fill.fore_color.rgb = RGBColor(254, 237, 238)
            else:
                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)
                
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx != 3 else PP_ALIGN.LEFT
            p.font.name = BODY_FONT
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_HIGHLIGHT if (row_idx == 2 and col_idx > 0) else COLOR_TEXT_MAIN
            if row_idx == 2:
                p.font.bold = True

    # ==========================================
    # SLIDE 6: 실증 성적 분석 데이터 (2024학년도 1학년 코호트 분석 결과) [신규!]
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_title(slide6, "통계로 입증된 성적 향상 레버리지와 관외 유입")
    add_header_and_logo(slide6, prs)
    
    content_box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf6 = content_box.text_frame
    tf6.word_wrap = True
    
    pts = [
        ("📈 중학 성적 vs 고교 석차의 높은 연계와 성공 가능성", 0, True),
        ("- 중학 석차백분율(%)과 고교 1학기말 석차 상관계수 0.8631로 매우 높은 신뢰도 입증", 1, False),
        ("🔥 꿈을 현실로 만드는 '내신 역전(도약)의 드라마' 선배 사례", 0, True),
        ("- 임지민 선배 (안골포중 출신): 고교 입학 26등 ➔ 1학년 1학기말 전교 1등 달성 (▲25계단 수직 상승!)", 1, False),
        ("- 김준민 선배 (진해남중 출신): 고교 입학 21등 ➔ 1학년 1학기말 전교 4등 달성 (▲17계단 도약!)", 1, False),
        ("- 윤영 선배 (진해남중 출신): 고교 입학 17등 ➔ 1학년 1학기말 전교 3등 달성 (▲14계단 도약!)", 1, False),
        ("🚗 관외/원거리 유입률 20.7% 및 기숙사(송학관/동백관) 완벽 적응", 0, True),
        ("- 신입생 중 창원 관내, 용원, 웅천 등 관외 출신 비율이 20.7% (63명) 달성", 1, False),
        ("- 기숙사 정주 여건을 기반으로 관외 출신 학생들도 최상위권(전교 1등 임지민 선배 등)에 다수 포진", 1, True)
    ]
    
    for idx, (text, level, highlight) in enumerate(pts):
        p = tf6.paragraphs[0] if idx == 0 else tf6.add_paragraph()
        p.text = text
        p.level = level
        p.line_spacing = 1.25
        p.space_after = Pt(6)
        p.font.name = BODY_FONT
        p.font.size = Pt(24 if level == 0 else 19)
        p.font.color.rgb = COLOR_TITLE if level == 0 or not highlight else COLOR_HIGHLIGHT
        if highlight:
            p.font.bold = True

    # ==========================================
    # SLIDE 7: 의·치·약학 계열 최상위 선배 대입 로드맵 [신규!]
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_title(slide7, "의·치·약학 메디컬 계열 진학을 향한 최상위권 선배 사례")
    add_header_and_logo(slide7, prs)
    
    # 3명의 최상위 선배 카드를 가로 배치
    card_width = Inches(3.6)
    card_height = Inches(4.2)
    card_y = Inches(2.2)
    card_xs = [Inches(0.8), Inches(4.8), Inches(8.8)]
    
    med_students = [
        ("유도훈 선배 (2학년 1반)", "• 내신 평균: 1.08 등급\n• 지망 학과: 연세대학교 의예과\n• 특징: 중학 백분율 1.67%로 입학, 철저한 내신 유지와 생명과학 의학 융합 탐구 세특 최우수 관리", COLOR_TITLE),
        ("임지민 선배 (2학년 2반)", "• 내신 평균: 1.20 등급\n• 지망 학과: 경상국립대 의예과\n• 특징: 중학 백분율 0.83%로 안골포중에서 진해고로 입학하여 지역인재 특별전형 수시 완벽 대비", COLOR_HIGHLIGHT),
        ("박준형 선배 (2학년 5반)", "• 내신 평균: 2.00 등급\n• 지망 학과: 경상국립대 의예과\n• 특징: 기숙사 자기주도적 학습 지원과 심화 탐구 공동체 활동으로 종합 전형 및 의예과 수시 공략", COLOR_TITLE)
    ]
    
    for i, (student, desc, color) in enumerate(med_students):
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_xs[i], card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        tb = slide7.shapes.add_textbox(card_xs[i] + Inches(0.1), card_y + Inches(0.1), card_width - Inches(0.2), card_height - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_s = tf.paragraphs[0]
        p_s.text = student
        p_s.font.name = TITLE_FONT
        p_s.font.size = Pt(20)
        p_s.font.bold = True
        p_s.font.color.rgb = color
        p_s.space_after = Pt(14)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = BODY_FONT
        p_d.font.size = Pt(15)
        p_d.line_spacing = 1.3
        p_d.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 8: 교내 설명회 - AI 기반 학생 지도 시스템 시연
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_title(slide8, "공교육을 혁신하는 AI 기반 첨단 학생 지도 시스템")
    add_header_and_logo(slide8, prs)
    
    # 3개의 시연 카드를 가로 배치
    card_width = Inches(3.6)
    card_height = Inches(4.2)
    card_y = Inches(2.2)
    card_xs = [Inches(0.8), Inches(4.8), Inches(8.8)]
    
    ai_tools = [
        ("① AI 국어 문법 Socratic 튜터", "학생이 오개념을 입력하면 AI가 소크라테스식 대화 유도로 스스로 개념을 귀납적 탐구하도록 이끄는 1:1 학습 도우미", COLOR_TITLE),
        ("② AI 세특 자동화 시스템", "구글 시트 백엔드와 Gemini API 연동을 통해 학생의 수행 이력을 분석하여 개성 넘치고 고유한 대입용 맞춤 세특 초안 자동 생성", COLOR_HIGHLIGHT),
        ("③ AI 학부모 알림 메시지 포털", "바쁜 담임 선생님이 적은 관찰 단서를 격식 있고 감동적인 소통 문구로 3초 만에 변환하여 학부모 안심 케어 서비스 구축", COLOR_TITLE)
    ]
    
    for i, (title, desc, color) in enumerate(ai_tools):
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_xs[i], card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        tb = slide8.shapes.add_textbox(card_xs[i] + Inches(0.1), card_y + Inches(0.1), card_width - Inches(0.2), card_height - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = TITLE_FONT
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.space_after = Pt(14)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = BODY_FONT
        p_d.font.size = Pt(16)
        p_d.line_spacing = 1.3
        p_d.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 9: 마무리 요약
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_title(slide9, "대입 성공을 위한 최고의 동반자: 진해고등학교")
    add_header_and_logo(slide9, prs)
    
    content_box = slide9.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf9 = content_box.text_frame
    tf9.word_wrap = True
    
    summary_points = [
        ("1. 내신 확보의 절대적 유리함 & 남학교 환경", " - 대규모 수강 과목 개설로 넓은 등급 문호 확보 및 공정한 내신 등급 경쟁", COLOR_TITLE),
        ("2. 교내 과목 개설 충족율 100% 보장", " - 소규모 폐강 걱정 없이 본인의 대입 진로에 필요한 모든 선택 과목 실제 수강", COLOR_HIGHLIGHT),
        ("3. 실증 통계로 검증된 성적 역전 드라마", " - 피어슨 상관계수 0.8631 및 입학 성적 대비 1학기말 성취 상승(임지민 선배 26등➔1등)", COLOR_TITLE),
        ("4. 의·치·약학 메디컬 진학 체계적 케어", " - 최상위 의예과 지망 선배(유도훈 1.08등급 등) 1:1 맞춤 관리 및 기숙사 최적 지원", COLOR_TITLE)
    ]
    
    for idx, (title, desc, color) in enumerate(summary_points):
        p_t = tf9.paragraphs[0] if idx == 0 else tf9.add_paragraph()
        p_t.text = title
        p_t.font.name = TITLE_FONT
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        p_d = tf9.add_paragraph()
        p_d.text = desc
        p_d.font.name = BODY_FONT
        p_d.font.size = Pt(17)
        p_d.font.color.rgb = COLOR_TEXT_MAIN
        p_d.space_after = Pt(8)
        
    prs.save("jinhae_admission_presentation.pptx")
    print("성공적으로 jinhae_admission_presentation.pptx 파일이 생성되었습니다.")

if __name__ == "__main__":
    create_admission_presentation()
