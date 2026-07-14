import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# 援ш? ?쒕씪?대툕 ?낅줈?쒖슜 ?ㅼ젙
FOLDER_ID = '1lxKMH8ssyeHSEv3rI2a24pD8qCT6FXo1'
SCOPES = ['https://www.googleapis.com/auth/drive']

def create_presentation():
    prs = Presentation()
    
    # ?덉씠?꾩썐 ?뺤쓽 (0: ?쒕ぉ, 1: ?쒕ぉ+?댁슜, 5: ?쒕ぉ留?
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # [1] ?쒕ぉ ?щ씪?대뱶
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "2027 ?섎뒫?밴컯 臾명븰: 媛덈옒蹂듯빀 05"
    subtitle.text = "?덈궃?ㅽ뿄 <洹쒖썝媛> & ?묒옄 誘몄긽 <?명쁽?뺥썑??\n?꾨꼍 ?대? 諛?媛쒕뀗 ?뺣━"
    
    # [2] 洹쒖썝媛 媛쒓?
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "1. ?덈궃?ㅽ뿄 <洹쒖썝媛> 媛쒓?"
    tf = slide.placeholders[1].text_frame
    tf.text = "媛덈옒: ?대갑 媛??洹쒕갑 媛??, ?묐컲 媛??
    p = tf.add_paragraph()
    p.text = "二쇱젣: 遊됯굔???ы쉶 ?쒕룄 ?꾨옒 ?낆닔怨듬갑?섎뒗 ?ъ꽦???쒗깂怨??먮쭩"
    p = tf.add_paragraph()
    p.text = "?섏쓽: ?꾩쟾 媛???ㅻ옒??洹쒕갑 媛??
    p = tf.add_paragraph()
    p.text = "?듭떖 媛먯긽 ?ъ씤??"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "怨쇨굅???꾨쫫?ㅼ?(?ㅻ퉰?붿븞)怨??꾩옱??珥덈씪??硫대ぉ媛利? ?鍮?
    p.level = 2
    p = tf.add_paragraph()
    p.text = "?⑦렪??????먮쭩怨?湲곕떎由쇱씠?쇰뒗 ?댁쨷???쒕룄"
    p.level = 2
    
    # [3] 洹쒖썝媛 以꾧굅由?(湲? ??
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "2. <洹쒖썝媛> 二쇱슂 ?댁슜 (湲???"
    tf = slide.placeholders[1].text_frame
    tf.text = "??湲?壅?: 怨쇨굅 ?뚯긽 諛??좎꽭 ?쒗깂"
    p = tf.add_paragraph(); p.text = "15, 16??利덉쓬???덇컳?????쇰?? 怨좎슫 ?쇨뎬(?ㅻ퉰?붿븞)??吏?붿쑝??; p.level = 1
    p = tf.add_paragraph(); p.text = "?꾩옱???덈Т ?숈뼱 諛됱궡?ㅻ윭??紐⑥뒿(硫대ぉ媛利????섏뿀?뚯쓣 ?쒗깂"; p.level = 1
    p = tf.add_paragraph(); p.text = "議곕Ъ(?좂돥)???쒓린瑜??볧븯硫?遺洹?ㅻ궓???붿옄 ?놁쓬???ы띁??; p.level = 1
    p = tf.add_paragraph(); p.text = "??????: 諛⑺깢???⑦렪 ?먮쭩"
    p = tf.add_paragraph(); p.text = "?μ븞 ?좏삊 寃쎈컯???명깢?섍퀬 寃쎈컯??臾대━??? ?댁슱由щ뒗 ?⑦렪"; p.level = 1
    p = tf.add_paragraph(); p.text = "??湲곗깮吏??쇱쑀?????앷꼈?붿? 諛깅쭏湲덊렪?쇰줈 ?섍컙 ???뚯떇???딄?"; p.level = 1
    
    # [4] 洹쒖썝媛 以꾧굅由?(?? 寃?
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "3. <洹쒖썝媛> 二쇱슂 ?댁슜 (??寃?"
    tf = slide.placeholders[1].text_frame
    tf.text = "????饔?: ?낆닔怨듬갑???좊떖??
    p = tf.add_paragraph(); p.text = "珥덇?吏묒뿉 ?⑥뼱吏??鍮꾩? ??吏???뚮━, ?ㅼ넄(洹?쒕씪誘? ?곕뒗 ?뚮━???쒕쫫 怨좎“"; p.level = 1
    p = tf.add_paragraph(); p.text = "嫄곕Ц怨??밴린湲?瑜??硫??몃줈????щ옒??援ш끝媛꾩옣???딆뼱吏?; p.level = 1
    p = tf.add_paragraph(); p.text = "??寃?永?: 湲곌뎄???섎챸 泥대뀗"
    p = tf.add_paragraph(); p.text = "?쎌닔(?μ븷臾? ?볦뿉 ?뚯떇議곗감 ?놁쓬"; p.level = 1
    p = tf.add_paragraph(); p.text = "李⑤씪由?二쎌뼱 ?ㅻ옖 ?ㅼ쓽 ???뷀몴泥쒕뀈??蹂꾪븰)???섏뼱 留뚮굹蹂쇨퉴 泥대뀗??; p.level = 1

    # [5] 洹쒖썝媛 媛쒕뀗
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "4. <洹쒖썝媛> ?듭떖 媛쒕뀗: 媛먯젙?댁엯 vs 媛앷????곴?臾?
    tf = slide.placeholders[1].text_frame
    tf.text = "??媛먯젙?댁엯 (Empathy)"
    p = tf.add_paragraph(); p.text = "??곸뿉 ?붿옄??媛먯젙???ъ쁺?섏뿬 ??곸씠 ?붿옄? ?묎컳???먮겮??寃껋쿂??臾섏궗"; p.level = 1
    p = tf.add_paragraph(); p.text = "?곸슜: ?ㅼ넄 (洹?쒕씪誘? - ?붿옄???ы뵒??'?곕뒗 洹?쒕씪誘?濡??ъ쁺??; p.level = 2; p.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
    
    p = tf.add_paragraph(); p.text = "??媛앷????곴?臾?(Objective Correlative)"
    p = tf.add_paragraph(); p.text = "媛먯젙???섍린?섎뒗 紐⑤뱺 援ъ껜???щЪ濡? 媛먯젙?댁엯源뚯? ?ы븿?섎뒗 ????媛쒕뀗"; p.level = 1
    p = tf.add_paragraph(); p.text = "?곸슜: ?밴린湲?(嫄곕Ц怨? - ?곕뒗 ??곸씠 ?꾨땶, ?ы뵒???щ옒湲??꾪빐 ?곗＜?섎뒗 ?꾧뎄/留ㅺ컻泥?; p.level = 2; p.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)

    # [6] ?명쁽?뺥썑??媛쒓?
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "5. ?묒옄 誘몄긽 <?명쁽?뺥썑?? 媛쒓? 諛?諛곌꼍"
    tf = slide.placeholders[1].text_frame
    tf.text = "媛덈옒: 援?Ц ?뚯꽕, ??궗 ?뚯꽕"
    p = tf.add_paragraph(); p.text = "二쇱젣: ?명쁽?뺥썑??怨좉껐???뺤꽦怨??쒕젴洹밸났, ?ы븘洹??
    p = tf.add_paragraph(); p.text = "諛곌꼍: 議곗꽑 ?꾧린 ?숈쥌 ?(湲곗궗?섍뎅, 媛묒닠?섍뎅)"
    p = tf.add_paragraph(); p.text = "?섏쓽: ?쒖쨷濡? 怨꾩텞?쇨린? ?붾텋??議곗꽑 ?쒕? 3? 沅곸쨷 臾명븰"
    p = tf.add_paragraph(); p.text = "?듭떖 媛먯긽 ?ъ씤??"
    p.level = 1
    p = tf.add_paragraph(); p.text = "?좉탳???뺣ぉ(?ш린 湲덉?, ?쒖쥌, ?몃궡)??蹂댁뿬二쇰뒗 ?덈??????????뺤긽??; p.level = 2

    # [7] ?명쁽?뺥썑??以꾧굅由?    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "6. <?명쁽?뺥썑?? ?꾩껜 以꾧굅由?留λ씫"
    tf = slide.placeholders[1].text_frame
    tf.text = "??諛쒕떒: 15?몄쓽 ?명쁽?뺥썑媛 怨꾨퉬媛 ?섎굹 ?꾨뱾???녹? 紐삵븿"
    p = tf.add_paragraph(); p.text = "???꾧컻: ?숈쥌???낃턿?쒗궓 ?ν씗鍮덉씠 ?꾨뱾(寃쎌쥌)???녠퀬 ?명쁽?뺥썑瑜??뚰빐?섍린 ?쒖옉??
    p = tf.add_paragraph(); p.text = "???꾧린(蹂몃Ц): ?ν씗鍮덉쓽 紐⑦븿?쇰줈 ?댁꽦???껋? ?숈쥌???명쁽?뺥썑瑜??쒖씤?쇰줈 ?먯쐞?쒗궡 (湲곗궗?섍뎅)"
    p = tf.add_paragraph(); p.text = "洹몃윭???명쁽?뺥썑???먮쭩 ?놁씠 ?좉탳???꾨━瑜?吏?ㅻŉ ?멸퀬??; p.level = 1
    p = tf.add_paragraph(); p.text = "???덉젙/寃곕쭚: ?ν씗鍮덉쓽 ?쒖뾾??援먮쭔怨??낇뻾"; p.level = 0
    p = tf.add_paragraph(); p.text = "?숈쥌??源⑤떕怨??명쁽?뺥썑瑜?蹂듭쐞?쒗궡(媛묒닠?섍뎅). ?댄썑 ?μ뵪???ъ빟??諛쏆븘 ?ы븘洹???ㅽ쁽"; p.level = 1

    # [8] ?명쁽?뺥썑???몃Ъ 媛덈벑
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "7. ?몃Ъ 媛??由쎄낵 ?좉탳???ㅻ━愿"
    tf = slide.placeholders[1].text_frame
    tf.text = "?뮕 ?꾪삎???몃Ъ???댁꽭??沅뚯꽑吏뺤븙 ?뱁샇"
    p = tf.add_paragraph(); p.text = "?명쁽?뺥썑 (?덈?????"
    p = tf.add_paragraph(); p.text = "?듭슱???먯쐞 泥섎텇?먮룄 ???븯吏 ?딆쓬"; p.level = 1
    p = tf.add_paragraph(); p.text = "援곗＜??????앹뾾??遺??惹?쓿)怨??쒖쥌"; p.level = 1
    p = tf.add_paragraph(); p.text = "?ν씗鍮?(?덈?????"; p.level = 0
    p = tf.add_paragraph(); p.text = "?꾪넻??媛遺?μ젣?먯꽌 媛??湲덇린?쒗븯??'?ш린(吏덊닾)'? 李몄냼???붿떊"; p.level = 1
    p = tf.add_paragraph(); p.text = "?숈쥌 (?덈???沅뚮젰??"; p.level = 0
    p = tf.add_paragraph(); p.text = "紐⑦븿???덉씠 硫吏留? 援?솗?대?濡?吏곸젒?곸씠怨?媛뺣젰??鍮꾪뙋?먯꽌??踰쀬뼱??; p.level = 1

    # [9] ?명쁽?뺥썑???뱀쭠
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "8. ?꾩닔 媛쒕뀗: ?쒖닠?먯쓽 媛쒖엯(?몄쭛?먯쟻 ?쇳룊)"
    tf = slide.placeholders[1].text_frame
    tf.text = "?뺤쓽: ?꾩????묎? ?쒖젏 ?꾨옒?먯꽌 ?쒖닠?먭? ?ш굔 以묎컙???섏꽌 ?몃Ъ?대굹 ?곹솴?????媛먯젙怨??됯?瑜?吏곸젒 ?쒖텧?섎뒗 寃?
    
    p = tf.add_paragraph(); p.text = "蹂몃Ц ?덉떆:"
    p = tf.add_paragraph(); p.text = '"?댁컡 媛?⑥튂 ?꾨땲?섎━?? 留뚮Ъ???ы띁?섍퀬 ?섎뒛??媛먮룞?덈떎 ?⑥씠 ?댁컡 鍮덈쭚?대━??"' ; p.level = 1
    
    p = tf.add_paragraph(); p.text = "?④낵:"
    p = tf.add_paragraph(); p.text = "?쒖닠?먭? ?명쁽?뺥썑瑜??몄븷?섏뿬"; p.level = 1
    p = tf.add_paragraph(); p.text = "?낆옄?ㅼ씠 ?쒖닠?먯쓽 吏숈? 媛먯젙???숉솕?섍쾶 留뚮벀"; p.level = 1
    p = tf.add_paragraph(); p.text = "?먯뿰?ㅻ젅 ?명쁽?뺥썑???숈젙?섍퀬 ?ν씗鍮덉쓣 鍮꾪뙋?섍쾶 ?좊룄"; p.level = 1

    prs.save("worksheet_presentation.pptx")
    print("Local presentation saved as worksheet_presentation.pptx")

def upload_presentation():
    creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    service = build('drive', 'v3', credentials=creds)
    
    file_metadata = {
        'name': '[?щ씪?대뱶] 2027 ?섎뒫?밴컯 臾명븰: 媛덈옒蹂듯빀 05',
        'parents': [FOLDER_ID],
        # ?닿쾬???듭떖! PPTX瑜?援ш? ?щ씪?대뱶濡?蹂??        'mimeType': 'application/vnd.google-apps.presentation'
    }
    
    media = MediaFileUpload('worksheet_presentation.pptx',
                            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                            resumable=True)
                            
    print("Uploading to Google Drive and converting to Slides...")
    file = service.files().create(body=file_metadata, media_body=media, fields='id, webViewLink').execute()
    print(f"Presentation uploaded! Link: {file.get('webViewLink')}")

if __name__ == "__main__":
    create_presentation()
    upload_presentation()
